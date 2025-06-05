import os
import json
import re
import logging
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
 
# Configure logging
logging.basicConfig(filename='resume_generator.log', level=logging.ERROR, format='%(asctime)s:%(levelname)s:%(message)s')
 
FONT="Helvetica"
 
def layout5(parsed_result, path, image_path=None):
    try:
        if not os.path.exists(path):
            os.makedirs(path)
        if isinstance(parsed_result, str):
            data = json.loads(parsed_result)
        elif isinstance(parsed_result, dict):
            data = parsed_result
        else:
            raise ValueError("parsed_result must be a JSON string or a dictionary")
        print(type(data))
        # Create a presentation object
        prs = Presentation()
 
 
        # Change slide size to 16:9 widescreen format (13.33 inches x 7.5 inches)
        prs.slide_width = Cm(67.729)
        prs.slide_height = Cm(38.1)
 
        # Add a slide with a title layout
        slide_layout = prs.slide_layouts[5]  # Use the blank layout
        slide = prs.slides.add_slide(slide_layout)
 
 
        # Define the position and size of the text boxes
        left0, top0, width0, height0 = Cm(2), Cm(2), Cm(14), Cm(6) #Name
        left1, top1, width1, height1 = Cm(2), Cm(5), Cm(31), Cm(7) #Summary
        left2, top2, width2, height2 = Cm(2), Cm(12), Cm(8), Cm(9) #Roles played
        left3, top3, width3, height3 = Cm(12), Cm(12), Cm(7), Cm(8)#Area's of Experties
        left4, top4, width4, height4 = Cm(24), Cm(12), Cm(7), Cm(8)#tools
        left5, top5, width5, height5 = Cm(12), Cm(20), Cm(8), Cm(6)#Industry sectors
        left6, top6, width6, height6 = Cm(2), Cm(28), Cm(6), Cm(6)#Consulting engagement
        left7, top7, width7, height7 = Cm(12), Cm(28), Cm(21), Cm(6)#Education
        left8, top8, width8, height8 = Cm(35), Cm(1.5), Cm(32), Cm(1.5)#Experience and Accomplishments
        left9, top9, width9, height9 = Cm(35), Cm(3), Cm(32), Cm(34)#Education
 
        # Add a shape to cover half of the slide with light gray background
        left_bg, top_bg, width_bg, height_bg = Cm(0), Cm(-0), Cm(33.86), Cm(38.1) # Half slide
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_bg, top_bg, width_bg, height_bg)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(242, 242, 242) # Light gray
 
        # Remove the border
        shape.line.fill.background()
       
        # Add image to the top right corner if provided
        if image_path and os.path.exists(image_path):
            img_left = Cm(25.63)
            img_top = Cm(0.48)
            img_width = Cm(5.99)
            img_height = Cm(5.99)
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)
 
        # Add the first text box
        text_box0 = slide.shapes.add_textbox(left0, top0, width0, height0)
        text_frame0 = text_box0.text_frame
        text_frame0.text =data.get("Name", "Unnamed")
        p1 = text_frame0.paragraphs[0]
        run1 = p1.runs[0]
        font1 = run1.font
        font1.name = FONT
        font1.size = Pt(30)
        font1.bold = True
        font1.color.rgb = RGBColor(0, 0, 139)
        contact_info = [
            data.get("Email", ""),
            data.get("Phone", ""),
            data.get("LinkedIn", ""),
            data.get("Address", "")
        ]
        for info in contact_info:
            p2=text_frame0.add_paragraph()
            p2.level = 0
            p2.text=info
            run2 = p2.add_run()
            font2 = run2.font
            font2.name = "Helvetica"
            font2.size = Pt(24)
            font2.color.rgb = RGBColor(0, 102, 204) # Blue
 
        # Add the first text box
        text_box1 = slide.shapes.add_textbox(left1, top1, width1, height1)
        text_frame1 = text_box1.text_frame
        text_frame1.word_wrap = True # Enable word wrap
        text_frame1.text = "Summary"
        p1 = text_frame1.paragraphs[0]
        run1 = p1.runs[0]
        font1 = run1.font
        font1.name = FONT
        font1.size = Pt(20)
        font1.bold = True
        font1.color.rgb = RGBColor(0, 0, 139)
        p2=text_frame1.add_paragraph()
        p2.text=data.get("Summary") #type: ignore
        p2.level = 0
        run2 = p2.add_run()
        font2 = run2.font
        font2.name = "Helvetica"
        font2.size = Pt(18)
        #font2.color.rgb = RGBColor(0, 102, 204) # Blue
 
        # Add the second text box
        text_box2 = slide.shapes.add_textbox(left2, top2, width2, height2)
        text_frame2 = text_box2.text_frame
        text_frame2.word_wrap = True
        text_frame2.text = "Roles played"
        p1 = text_frame2.paragraphs[0]
        run1 = p1.runs[0]
        font1 = run1.font
        font1.name = FONT
        font1.size = Pt(20)
        font1.bold = True
        font1.color.rgb = RGBColor(0, 0, 139)
           
        roles_played = data.get("Roles Played") or []
        if isinstance(roles_played, str):
            roles_played = [roles_played]
        for role in roles_played:
            p2 = text_frame2.add_paragraph()
            p2.text = f"- {role}"
            p2.level = 0
            run2 = p2.runs[0]
            font2 = run2.font
            font2.name = FONT
            font2.size = Pt(18)
            #font2.color.rgb = RGBColor(0, 102, 204)
 
        # Add the third text box
        text_box3 = slide.shapes.add_textbox(left3, top3, width3, height3)
        text_frame3 = text_box3.text_frame
        text_frame3.text = "Areas of Expertise"
        p1 = text_frame3.paragraphs[0]
        run1 = p1.runs[0]
        font1 = run1.font
        font1.name = FONT
        font1.size = Pt(20)
        font1.bold = True
        font1.color.rgb = RGBColor(0, 0, 139)
       
        expertise_areas = data.get("Areas of Expertise") or []
        if isinstance(expertise_areas, str):
            expertise_areas = [expertise_areas]
        for area in expertise_areas:
            p2 = text_frame3.add_paragraph()
            p2.text = f"- {area}"
            p2.level = 0
            run2 = p2.runs[0]
            font2 = run2.font
            font2.name = FONT
            font2.size = Pt(18)
            #font2.color.rgb = RGBColor(0, 102, 204)
 
        # Add the fourth text box
        text_box4 = slide.shapes.add_textbox(left4, top4, width4, height4)
        text_frame4 = text_box4.text_frame
        text_frame4.word_wrap = True
        text_frame4.text = "Skills"
        p1 = text_frame4.paragraphs[0]
        run1 = p1.runs[0]
        font1 = run1.font
        font1.name = FONT
        font1.size = Pt(20)
        font1.bold = True
        font1.color.rgb = RGBColor(0, 0, 139)
       
        skills = data.get("Skills") or []
        if isinstance(skills, str):
            skills = [skills]
        elif skills is None:
            skills = []
           
        for skill in skills:
            if skill:
                p2 = text_frame4.add_paragraph()
                p2.text = f"- {skill}"
                p2.level = 0
                run2 = p2.runs[0]
                font2 = run2.font
                font2.name = FONT
                font2.size = Pt(18)
                #font2.color.rgb = RGBColor(0, 102, 204) # Blue
 
        # Add the fifth text box
        text_box5 = slide.shapes.add_textbox(left5, top5, width5, height5)
        text_frame5 = text_box5.text_frame
        text_frame5.text = "Industry Sectors"
        p1 = text_frame5.paragraphs[0]
        run1 = p1.runs[0]
        font1 = run1.font
        font1.name = FONT
        font1.size = Pt(20)
        font1.bold = True
        font1.color.rgb = RGBColor(0, 0, 139)
       
        sectors = data.get("Industry Sectors") or []
        if isinstance(sectors, str):
            sectors = [sectors]
        for sector in sectors:
            p2 = text_frame5.add_paragraph()
            p2.text = f"- {sector}"
            p2.level = 0
            run2 = p2.runs[0]
            font2 = run2.font
            font2.name = FONT
            font2.size = Pt(18)
            #font2.color.rgb = RGBColor(0, 102, 204)
 
        # Add the sixth text box
        text_box6 = slide.shapes.add_textbox(left6, top6, width6, height6)
        text_frame6 = text_box6.text_frame
        text_frame6.text = "Consulting Engagements"
        p1 = text_frame6.paragraphs[0]
        run1 = p1.runs[0]
        font1 = run1.font
        font1.name = FONT
        font1.size = Pt(20)
        font1.bold = True
        font1.color.rgb = RGBColor(0, 0, 139)
       
        consulting = data.get("Consulting Engagements", "N/A") or []
        if isinstance(consulting, str):
            consulting = [consulting]
        for item in consulting:
            p2 = text_frame6.add_paragraph()
            p2.text = f"- {item}"
            p2.level = 0
            run2 = p2.runs[0]
            font2 = run2.font
            font2.name = FONT
            font2.size = Pt(18)
            #font2.color.rgb = RGBColor(0, 102, 204) # Blue
 
        # Add the seventh text box
        text_box7 = slide.shapes.add_textbox(left7, top7, width7, height7)
        text_frame7 = text_box7.text_frame
        text_frame7.text = "Education & Certifications"
        text_frame7.word_wrap = True
        p1 = text_frame7.paragraphs[0]
        run1 = p1.runs[0]
        font1 = run1.font
        font1.name = FONT
        font1.size = Pt(20)
        font1.bold = True
        font1.color.rgb = RGBColor(0, 0, 139)
 
        education = data.get("Education or Academic Profile and Certifications", "N/A") or []
        if isinstance(education, dict):
            education = [education]
        elif isinstance(education, str):
            education = [education]
 
        for entry in education:
            if isinstance(entry, dict):
                degree = entry.get("Degree", "N/A")
                institution = entry.get("Institution", "N/A")
                duration = entry.get("Duration or Year", "")
                p2 = text_frame7.add_paragraph()
                p2.text = f"- {degree}, {institution} ({duration})"
            else:
                p2 = text_frame7.add_paragraph()
                p2.text = f"- {entry}"
            p2.level = 0
            run2 = p2.runs[0]
            font2 = run2.font
            font2.name = FONT
            font2.size = Pt(18)
            #font2.color.rgb = RGBColor(0, 102, 204) # Blue
           
        # Add and style the eighth text box
        text_box8 = slide.shapes.add_textbox(left8, top8, width8, height8)
        text_frame8 = text_box8.text_frame
        text_frame8.text = "Experience and Accomplishments"
        fill = text_box8.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(173, 216, 230) # Light blue
        p = text_frame8.paragraphs[0]
        run = p.runs[0]
        font = run.font
        font.name = FONT
        font.size = Pt(25)
        font.bold = True
        font.italic = False
        font.color.rgb = RGBColor(0, 0, 139) # Dark Blue
 
        # Add the ninth text box (Content)
        text_box9 = slide.shapes.add_textbox(left9, top9, width9, height9)
        text_frame9 = text_box9.text_frame
        text_frame9.word_wrap = True
 
        experience = data.get("Experience and Accomplishments")
        if not experience:
            experience = []
 
        if isinstance(experience, dict):
            experience = [experience]
        elif isinstance(experience, str):
            experience = [experience]
 
        for exp in experience:
            if isinstance(exp, dict):
                title = exp.get("Title", "")
                company = exp.get("Company", "")
                duration = exp.get("Duration", "N/A")
                responsibilities = exp.get("Detailed Roles and Responsibilities") or []
               
                # Add job title, company, and duration
                p1 = text_frame9.add_paragraph()
                p1.text = f"{title} at {company} ({duration})"
                p1.level = 0
                run1 = p1.runs[0]
                font1 = run1.font
                font1.name = FONT
                font1.size = Pt(20)
                font1.bold = True
                font1.color.rgb = RGBColor(0, 102, 204)
 
                # Add responsibilities
                if isinstance(responsibilities, str):
                    responsibilities = [responsibilities]
                for resp in responsibilities:
                    if resp:
                        p2 = text_frame9.add_paragraph()
                        p2.text = f"- {resp}"
                        p2.level = 1
                        run2 = p2.runs[0]
                        font2 = run2.font
                        font2.name = FONT
                        font2.size = Pt(18)
                        #font2.color.rgb = RGBColor(0, 102, 204)
            elif isinstance(exp, str):
                p = text_frame9.add_paragraph()
                p.text = f"- {exp}"
                p.level = 0
                run = p.runs[0]
                font = run.font
                font.name = FONT
                font.size = Pt(18)
                #font.color.rgb = RGBColor(0, 102, 204)
 
        # Save
        safe_name = re.sub(r'[\\/*?:"<>|]', "", data.get('Name', 'Unnamed'))
        filename = f"formatted_resume-{safe_name}.pptx"
        full_path = os.path.join(path, filename)
        prs.save(full_path)
        return full_path
        # prs.save('presentation.pptx')
 
        # print("Presentation created successfully with two text boxes.")
   
    except Exception as e:
        logging.error(f"An error occurred while generating the resume: {e}")
        return None